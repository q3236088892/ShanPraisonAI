"""PPT 自动生成工作流 — 本地工具集。

提供 PPT 代码模板读取、python-pptx 脚本执行、生成结果验证三个工具。
这些工具在本地执行，不消耗 LLM token。
"""

from __future__ import annotations

import json
import os
import subprocess
import sys
import tempfile


def read_template(style: str = "dark") -> str:
    """读取预制的 python-pptx 代码模板。

    Args:
        style: 风格名称。dark = 深色商务风，light = 浅色简约风。

    Returns:
        包含模板代码的字符串，其中 ``# SLIDES_PLACEHOLDER`` 标记需被替换为实际幻灯片调用。
    """
    templates = {
        "dark": _TEMPLATE_DARK,
        "light": _TEMPLATE_LIGHT,
    }
    return templates.get(style, _TEMPLATE_DARK)


def generate_pptx(python_code: str, output_path: str = "output.pptx") -> str:
    """执行 python-pptx 脚本生成 PPT 文件。

    Args:
        python_code: 完整的 python-pptx Python 脚本内容。
                     脚本中的 ``OUTPUT_PATH`` 会被替换为 *output_path*。
        output_path: 输出文件路径。

    Returns:
        "SUCCESS: <path> (<size> bytes)" 或 "ERROR: <message>"。
    """
    script = python_code.replace("OUTPUT_PATH", output_path)

    # 写临时脚本
    tmp_fd, tmp_script = tempfile.mkstemp(suffix=".py", prefix="_ppt_gen_")
    try:
        with os.fdopen(tmp_fd, "w", encoding="utf-8") as f:
            f.write(script)

        result = subprocess.run(
            [sys.executable, tmp_script],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            err = (result.stderr or result.stdout or "unknown error")[:800]
            return f"ERROR: {err}"

        if os.path.exists(output_path):
            size = os.path.getsize(output_path)
            return f"SUCCESS: {output_path} ({size} bytes)"

        return "ERROR: script finished but output file not found"
    except subprocess.TimeoutExpired:
        return "ERROR: script execution timed out (120s)"
    except Exception as exc:
        return f"ERROR: {exc}"
    finally:
        if os.path.exists(tmp_script):
            os.remove(tmp_script)


def verify_pptx(file_path: str) -> str:
    """验证生成的 PPTX 文件：存在性、大小、页数。

    Args:
        file_path: PPTX 文件路径。

    Returns:
        JSON 字符串，包含 status (PASS/FAIL)、slide_count、size_bytes 等字段。
    """
    if not os.path.exists(file_path):
        return json.dumps({"status": "FAIL", "reason": "file not found"}, ensure_ascii=False)

    size = os.path.getsize(file_path)
    if size < 1000:
        return json.dumps(
            {"status": "FAIL", "reason": f"file too small: {size} bytes"},
            ensure_ascii=False,
        )

    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        # 提取每页标题用于报告
        slide_titles = []
        for slide in prs.slides:
            title_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    title_text = shape.text_frame.text.strip()
                    if title_text:
                        break
            slide_titles.append(title_text or "(untitled)")

        return json.dumps(
            {
                "status": "PASS",
                "file": file_path,
                "size_bytes": size,
                "slide_count": slide_count,
                "slide_titles": slide_titles,
            },
            ensure_ascii=False,
        )
    except ImportError:
        return json.dumps(
            {"status": "FAIL", "reason": "python-pptx not installed"},
            ensure_ascii=False,
        )
    except Exception as exc:
        return json.dumps(
            {"status": "FAIL", "reason": str(exc)},
            ensure_ascii=False,
        )


# ─── 预制模板 ────────────────────────────────────────────────

_TEMPLATE_DARK = r'''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色方案（深色商务风）──
BG_TITLE   = RGBColor(0x1a, 0x1a, 0x2e)
BG_CONTENT = RGBColor(0x0f, 0x0f, 0x23)
CLR_TITLE  = RGBColor(0xff, 0xff, 0xff)
CLR_HEADING = RGBColor(0x64, 0xb5, 0xf6)
CLR_BODY   = RGBColor(0xe0, 0xe0, 0xe0)
CLR_SUB    = RGBColor(0xaa, 0xaa, 0xcc)
CLR_ACCENT = RGBColor(0x42, 0xa5, 0xf5)

def _set_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_title_slide(title, subtitle=""):
    """添加封面/结尾页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_TITLE)
    # 装饰线
    from pptx.util import Emu
    line = slide.shapes.add_shape(1, Inches(3), Inches(2.3), Inches(7), Emu(36000))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CLR_TITLE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = CLR_SUB
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets):
    """添加内容页。bullets 为字符串列表。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CONTENT)
    # 标题栏底色
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x30)
    title_bar.line.fill.background()
    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CLR_HEADING
    # 要点
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u25b8  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(14)

def add_two_column_slide(title, left_bullets, right_bullets):
    """添加双栏内容页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CONTENT)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CLR_HEADING
    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.4))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = f"\u25b8  {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(10)
    # 右栏
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(5.4))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.text = f"\u25b8  {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(10)

# SLIDES_PLACEHOLDER

prs.save("OUTPUT_PATH")
'''

_TEMPLATE_LIGHT = r'''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色方案（浅色简约风）──
BG_TITLE   = RGBColor(0xff, 0xff, 0xff)
BG_CONTENT = RGBColor(0xf8, 0xf9, 0xfa)
CLR_TITLE  = RGBColor(0x21, 0x21, 0x21)
CLR_HEADING = RGBColor(0x1a, 0x73, 0xe8)
CLR_BODY   = RGBColor(0x42, 0x42, 0x42)
CLR_SUB    = RGBColor(0x75, 0x75, 0x75)
CLR_ACCENT = RGBColor(0x1a, 0x73, 0xe8)

def _set_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_title_slide(title, subtitle=""):
    """添加封面/结尾页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_TITLE)
    from pptx.util import Emu
    line = slide.shapes.add_shape(1, Inches(4), Inches(2.3), Inches(5.333), Emu(28000))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CLR_TITLE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = CLR_SUB
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets):
    """添加内容页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CONTENT)
    # 顶部装饰线
    from pptx.util import Emu
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Emu(36000))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CLR_HEADING
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(14)

def add_two_column_slide(title, left_bullets, right_bullets):
    """添加双栏内容页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CONTENT)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CLR_HEADING
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.4))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(10)
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(5.4))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(10)

# SLIDES_PLACEHOLDER

prs.save("OUTPUT_PATH")
'''
